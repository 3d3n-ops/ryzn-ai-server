# main.py
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, WebSocket, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
import os
import shutil
import uuid
import logging
from typing import Optional, List, Dict, Any
from pydantic import BaseModel
from langchain_groq import ChatGroq
from gtts import gTTS
from dotenv import load_dotenv
from pathlib import Path
import re
from pydub import AudioSegment
import time
import json
import hashlib
import threading
from pdfminer.high_level import extract_text as extract_pdf_text
from queue import Queue
from threading import Thread
import asyncio
from fastapi import WebSocketDisconnect
import base64
import whisper
import tempfile
from .config import Settings, get_settings
from starlette.exceptions import HTTPException as StarletteHTTPException
from starlette.responses import JSONResponse
from fastapi.middleware.gzip import GZipMiddleware
# Add imports for PowerPoint and Word documents
from pptx import Presentation
import docx
import io

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Load environment variables first
load_dotenv()

# Get settings
settings = get_settings()

# Set up caching
CACHE_DIR = Path("cache")
CACHE_DIR.mkdir(exist_ok=True)

def get_cache_key(file_content: str, output_type: str) -> str:
    """Generate a unique cache key based on file content and output type."""
    content_hash = hashlib.md5(file_content.encode()).hexdigest()
    return f"{content_hash}_{output_type}"

def get_cached_response(cache_key: str) -> Optional[Dict[str, Any]]:
    """Retrieve cached response if it exists."""
    cache_file = CACHE_DIR / f"{cache_key}.json"
    if cache_file.exists():
        try:
            with open(cache_file, 'r') as f:
                cached_data = json.load(f)
                logger.debug(f"Cache hit for key: {cache_key}")
                return cached_data
        except Exception as e:
            logger.error(f"Error reading cache: {str(e)}")
    return None

def cache_response(cache_key: str, response_data: Dict[str, Any]) -> None:
    """Cache the response data."""
    try:
        cache_file = CACHE_DIR / f"{cache_key}.json"
        with open(cache_file, 'w') as f:
            json.dump(response_data, f)
        logger.debug(f"Cached response for key: {cache_key}")
    except Exception as e:
        logger.error(f"Error caching response: {str(e)}")

# Set up paths
BASE_DIR = Path(__file__).resolve().parent.parent.parent.parent
CREDENTIALS_DIR = BASE_DIR / "credentials"
CREDENTIALS_FILE = CREDENTIALS_DIR / "google-credentials.json"

# Ensure directories exist
CREDENTIALS_DIR.mkdir(exist_ok=True)
os.makedirs("uploads", exist_ok=True)
os.makedirs("static", exist_ok=True)

# Verify Google Cloud credentials (only for text-to-speech)
google_creds_path = os.getenv('GOOGLE_APPLICATION_CREDENTIALS')
if not google_creds_path:
    logger.error("GOOGLE_APPLICATION_CREDENTIALS not set in .env file")
    raise RuntimeError("GOOGLE_APPLICATION_CREDENTIALS not set in .env file. Please add it to your .env file.")

# Convert relative path to absolute if necessary
if not os.path.isabs(google_creds_path):
    google_creds_path = str(BASE_DIR / google_creds_path)

# Only check for credentials file in development mode
if os.environ.get('RENDER') != 'true' and not os.path.exists(google_creds_path):
    logger.error(f"Google Cloud credentials file not found at: {google_creds_path}")
    logger.error(f"Please place your google-credentials.json file in: {CREDENTIALS_DIR}")
    raise RuntimeError(f"Google Cloud credentials file not found. Please place your google-credentials.json file in: {CREDENTIALS_DIR}")

# Set the environment variable with the absolute path
os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = google_creds_path
logger.debug(f"Using Google Cloud credentials from: {google_creds_path}")

app = FastAPI()

# Add GZip middleware for compression
app.add_middleware(GZipMiddleware, minimum_size=1000)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=settings.ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Mount static directory for serving audio files
app.mount("/static", StaticFiles(directory="static"), name="static")

# Initialize LLM with Groq
groq_api_key = os.getenv("GROQ_API_KEY")
if not groq_api_key:
    raise ValueError("GROQ_API_KEY environment variable is not set")

llm = ChatGroq(
    model_name="llama-3.3-70b-versatile",  # Using llama model from Groq; mixtral model deprecated
    temperature=0.7,
    api_key=groq_api_key,
    streaming=False
)

# Add debug logging for API key
logger.debug(f"Using Groq API key: {groq_api_key[:8]}...")  # Only log first 8 chars for security

class ChatResponse(BaseModel):
    response: str
    user_query: str

class ChatRequest(BaseModel):
    content: dict
    user_query: str

class SummaryResponse(BaseModel):
    summary: str
    audio_url: Optional[str] = None
    notes: Optional[str] = None
    tts_audio_url: Optional[str] = None
    transcript: Optional[str] = None
    quiz: Optional[Dict[str, Any]] = None

# Add a simple rate limiter
class RateLimiter:
    def __init__(self, max_calls=5, time_period=60):
        self.max_calls = max_calls
        self.time_period = time_period
        self.calls = []
        self.lock = False
        self.last_reset = time.time()
        self.backoff_time = 5  # Initial backoff time in seconds
        self.max_backoff = 300  # Maximum backoff time (5 minutes)
    
    def can_call(self):
        current_time = time.time()
        
        # Reset calls if we're past the time period
        if current_time - self.last_reset >= self.time_period:
            self.calls = []
            self.last_reset = current_time
            self.backoff_time = 5  # Reset backoff time
            return True
        
        # Remove old calls
        self.calls = [call_time for call_time in self.calls if current_time - call_time < self.time_period]
        
        # Check if we can make a new call
        if len(self.calls) < self.max_calls and not self.lock:
            self.calls.append(current_time)
            return True
            
        # If we're rate limited, increase backoff time exponentially
        if not self.lock:
            self.backoff_time = min(self.backoff_time * 2, self.max_backoff)
            logger.warning(f"Rate limit reached. Waiting {self.backoff_time} seconds before retry.")
            time.sleep(self.backoff_time)
            return self.can_call()  # Try again after backoff
            
        return False
    
    def add_cooldown(self, seconds=10):
        """Add a cooldown period where no calls are allowed"""
        self.lock = True
        self.backoff_time = seconds
        threading.Timer(seconds, self._release_lock).start()
    
    def _release_lock(self):
        self.lock = False
        logger.debug("Rate limit cooldown period ended")

# Create rate limiter for Groq API with more generous limits
groq_limiter = RateLimiter(max_calls=30, time_period=60)  # 30 calls per minute

# Add a queue for processing chunks
class ChunkProcessor:
    def __init__(self, max_workers=3):
        self.queue = Queue()
        self.workers = []
        self.results = {}
        self.lock = threading.Lock()
        
        # Start worker threads
        for _ in range(max_workers):
            worker = Thread(target=self._process_chunks, daemon=True)
            worker.start()
            self.workers.append(worker)
    
    def add_chunk(self, chunk_id: str, chunk: str, callback):
        """Add a chunk to the processing queue"""
        self.queue.put((chunk_id, chunk, callback))
    
    def _process_chunks(self):
        """Worker thread to process chunks"""
        while True:
            try:
                chunk_id, chunk, callback = self.queue.get()
                
                # Process the chunk
                try:
                    result = callback(chunk)
                    with self.lock:
                        self.results[chunk_id] = result
                except Exception as e:
                    logger.error(f"Error processing chunk {chunk_id}: {str(e)}")
                    with self.lock:
                        self.results[chunk_id] = None
                
                self.queue.task_done()
            except Exception as e:
                logger.error(f"Error in chunk processor: {str(e)}")
    
    def get_result(self, chunk_id: str):
        """Get the result for a chunk"""
        with self.lock:
            return self.results.get(chunk_id)
    
    def wait_for_all(self, timeout=30):
        """Wait for all chunks to be processed with timeout"""
        start_time = time.time()
        while not self.queue.empty():
            if time.time() - start_time > timeout:
                raise TimeoutError("Summary generation timed out")
            time.sleep(0.1)
        return True

# Create a global chunk processor
chunk_processor = ChunkProcessor()

def process_chunk_with_llm(chunk: str) -> str:
    """Process a single chunk with the LLM"""
    prompt = f"""Please provide a concise summary of the following text:

    {chunk}

    Focus on the main points and key takeaways. Be direct and concise."""
    
    response = llm.invoke(prompt)
    return response.content if hasattr(response, 'content') else str(response)

def chunk_text(text: str, max_chunk_size: int = 4000) -> List[str]:
    """Split text into chunks of approximately equal size."""
    chunks = []
    current_chunk = []
    current_size = 0
    
    # Split text into sentences
    sentences = re.split('(?<=[.!?])\s+', text)
    
    for sentence in sentences:
        sentence_size = len(sentence)
        
        if current_size + sentence_size > max_chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_size = sentence_size
        else:
            current_chunk.append(sentence)
            current_size += sentence_size
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def generate_summary(text):
    try:
        # Add timeout for the entire operation
        start_time = time.time()
        timeout = 600  # 10 minutes timeout
        
        chunks = chunk_text(text)
        chunk_ids = [f"chunk_{i}" for i in range(len(chunks))]
        
        # Check cache for summaries
        cache_key = get_cache_key(text, "summary")
        cached_summary = get_cached_response(cache_key)
        if cached_summary:
            logger.debug("Using cached summary")
            return cached_summary.get("summary", "")
        
        # Process chunks in parallel with timeout
        for chunk_id, chunk in zip(chunk_ids, chunks):
            if time.time() - start_time > timeout:
                logger.error("Summary generation timed out")
                raise TimeoutError("Summary generation timed out")
                
            chunk_processor.add_chunk(chunk_id, chunk, process_chunk_with_llm)
        
        # Wait for all chunks to be processed with timeout
        chunk_timeout = 60  # 1 minute timeout per chunk
        while not chunk_processor.wait_for_all(timeout=chunk_timeout):
            if time.time() - start_time > timeout:
                logger.error("Summary generation timed out while waiting for chunks")
                raise TimeoutError("Summary generation timed out")
        
        # Collect results
        summaries = []
        for chunk_id in chunk_ids:
            result = chunk_processor.get_result(chunk_id)
            if result:
                summaries.append(result)
            else:
                logger.error(f"No result for chunk {chunk_id}")
        
        if not summaries:
            logger.error("No summaries were generated")
            raise HTTPException(status_code=500, detail="Failed to generate any summaries. Please try again.")
        
        # Combine summaries if needed
        if len(summaries) > 1:
            try:
                final_prompt = f"""Please provide a concise final summary combining these summaries:

                {' '.join(summaries)}

                Focus on the main points and key takeaways. Be direct and concise."""
                
                final_response = llm.invoke(final_prompt)
                final_summary = final_response.content if hasattr(final_response, 'content') else str(final_response)
            except Exception as e:
                logger.error(f"Error combining summaries: {str(e)}")
                final_summary = summaries[0]
        else:
            final_summary = summaries[0]
        
        # Cache the result
        cache_response(cache_key, {"summary": final_summary})
        return final_summary
            
    except TimeoutError as e:
        logger.error(f"Timeout generating summary: {str(e)}")
        raise HTTPException(status_code=504, detail="Summary generation timed out. Please try again with a smaller file or split it into parts.")
    except Exception as e:
        logger.error(f"Error generating summary: {str(e)}")
        if isinstance(e, HTTPException):
            raise e
        raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")

def generate_quiz(text):
    # Check cache first for this specific text and quiz
    cache_key = get_cache_key(text, "quiz")
    cached_quiz = get_cached_response(cache_key)
    if cached_quiz:
        logger.debug(f"Using cached quiz")
        return cached_quiz.get("quiz")
    
    # Truncate text to avoid token limits
    truncated_text = text[:4000]
    prompt = f"""With the context of this text: {truncated_text}, generate a quiz with 5-10 questions.
    Format your response as a JSON object with the following structure:
    {{
        "title": "Quiz Title",
        "description": "Brief description of the quiz",
        "questions": [
            {{
                "id": "q1",
                "question": "Question text",
                "options": ["Option A", "Option B", "Option C", "Option D"],
                "correctAnswer": "Correct option text",
                "explanation": "Explanation of why this is correct"
            }}
        ]
    }}
    Make sure to include 5-10 questions and format the response as valid JSON."""
    
    try:
        logger.debug("Generating quiz")
        
        # Check if we're within rate limits
        if not groq_limiter.can_call():
            logger.warning("Rate limit reached, waiting before generating quiz")
            time.sleep(20)
        
        response = llm.invoke(prompt)
        
        # Extract the JSON string from the response
        content = response.content
        # Find the JSON object in the response
        json_start = content.find('{')
        json_end = content.rfind('}') + 1
        if json_start == -1 or json_end == 0:
            raise ValueError("No valid JSON found in response")
            
        json_str = content[json_start:json_end]
        
        # Parse the JSON string into a Python dictionary
        quiz_data = json.loads(json_str)
        
        # Validate the structure
        required_fields = ['title', 'description', 'questions']
        for field in required_fields:
            if field not in quiz_data:
                raise ValueError(f"Missing required field: {field}")
                
        if not isinstance(quiz_data['questions'], list):
            raise ValueError("Questions must be a list")
            
        for question in quiz_data['questions']:
            required_question_fields = ['id', 'question', 'options', 'correctAnswer', 'explanation']
            for field in required_question_fields:
                if field not in question:
                    raise ValueError(f"Missing required field in question: {field}")
        
        # Cache the quiz data
        cache_response(cache_key, {"quiz": quiz_data})
        
        return quiz_data
        
    except Exception as e:
        logger.error(f"Error generating quiz: {str(e)}")
        
        # Check for rate limit errors
        error_msg = str(e)
        if "429" in error_msg or "Too Many Requests" in error_msg:
            logger.warning("Rate limit hit during quiz generation, adding cooldown")
            groq_limiter.add_cooldown(30)
            
            # Try to create a simple fallback quiz to avoid complete failure
            fallback_quiz = {
                "title": "Fallback Quiz",
                "description": "Sorry, we couldn't generate a detailed quiz due to API limits. Here's a simple quiz instead.",
                "questions": [
                    {
                        "id": "q1",
                        "question": "What's the primary topic of the uploaded document?",
                        "options": ["Please review the document", "Option B", "Option C", "Option D"],
                        "correctAnswer": "Please review the document",
                        "explanation": "This is a fallback quiz due to API limitations."
                    }
                ]
            }
            return fallback_quiz
            
        raise HTTPException(status_code=500, detail=f"Error generating quiz: {str(e)}")

def generate_notes(text):
    # Use LLM to generate notes
    prompt = f"Generate notes from the following text: {text}. Use markdown format. Use the following format: # Title\n\n## Summary\n\n## Key Points\n\n## Additional Information\n\n## References"
    try:
        logger.debug("Sending request to OpenAI")
        response = llm.invoke(prompt)
        return response.content
    except Exception as e:
        logger.error(f"Error generating notes: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error generating notes: {str(e)}")

def split_text_into_chunks(text: str, max_bytes: int = 4800) -> List[str]:
    """Split text into chunks that are within the byte limit."""
    chunks = []
    current_chunk = ""
    
    # Split by sentences (simple approach)
    sentences = re.split('(?<=[.!?])\s+', text)
    
    for sentence in sentences:
        # Check if adding this sentence would exceed the limit
        if len((current_chunk + sentence).encode('utf-8')) > max_bytes:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = sentence
        else:
            current_chunk = (current_chunk + " " + sentence).strip()
    
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks

def synthesize_long_audio(text: str, output_path: str) -> str:
    """Handle long text TTS by splitting into chunks and combining audio files."""
    # Initialize temp_files at the start to ensure it's always available
    temp_files = []
    try:
        client = texttospeech.TextToSpeechClient()
        chunks = split_text_into_chunks(text)
        combined_audio = None

        logger.debug(f"Split text into {len(chunks)} chunks")

        for i, chunk in enumerate(chunks):
            try:
                # Generate temporary file name
                temp_file = f"static/temp_{uuid.uuid4()}.mp3"
                temp_files.append(temp_file)
                
                synthesis_input = texttospeech.SynthesisInput(text=chunk)
                voice = texttospeech.VoiceSelectionParams(
                    language_code="en-US",
                    name="en-US-Studio-O",
                )
                audio_config = texttospeech.AudioConfig(
                    audio_encoding=texttospeech.AudioEncoding.MP3,
                    speaking_rate=1.0,
                    pitch=0.0,
                    sample_rate_hertz=24000
                )
                
                logger.debug(f"Synthesizing chunk {i+1}/{len(chunks)}")
                response = client.synthesize_speech(
                    input=synthesis_input,
                    voice=voice,
                    audio_config=audio_config
                )
                
                # Save temporary chunk
                with open(temp_file, "wb") as out:
                    out.write(response.audio_content)
                
                # Verify the temporary file was created
                if not os.path.exists(temp_file):
                    raise RuntimeError(f"Failed to create temporary audio file: {temp_file}")
                
                # Combine audio files using pydub
                if combined_audio is None:
                    combined_audio = AudioSegment.from_mp3(temp_file)
                else:
                    chunk_audio = AudioSegment.from_mp3(temp_file)
                    combined_audio += chunk_audio
                
                logger.debug(f"Successfully processed chunk {i+1}")

            except Exception as e:
                logger.error(f"Error processing chunk {i+1}: {str(e)}")
                raise RuntimeError(f"Error processing audio chunk {i+1}: {str(e)}")

        # Export final combined audio
        if combined_audio:
            logger.debug(f"Exporting final audio to {output_path}")
            combined_audio.export(output_path, format="mp3")
            
            if not os.path.exists(output_path):
                raise RuntimeError(f"Failed to create final audio file: {output_path}")
        else:
            raise RuntimeError("No audio was generated")
        
        return output_path
    except Exception as e:
        logger.error(f"Error in synthesize_long_audio: {str(e)}", exc_info=True)
        raise RuntimeError(f"Error generating audio: {str(e)}")
    finally:
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                    logger.debug(f"Cleaned up temporary file: {temp_file}")
            except Exception as e:
                logger.error(f"Error cleaning up temporary file {temp_file}: {str(e)}")

def transcribe_audio_file(audio_path):
    # Use Whisper to transcribe the audio file
    try:
        # Load the Whisper model
        model = whisper.load_model("base")
        result = model.transcribe(audio_path)
        return result["text"]
    except Exception as e:
        logger.error(f"Error transcribing audio: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error transcribing audio: {str(e)}")

@app.get("/api/status/{file_id}")
async def get_status(file_id: str):
    return {
        "status": chunk_processor.get_status(file_id),
        "progress": chunk_processor.get_progress(file_id),
        "estimated_time_remaining": chunk_processor.get_estimated_time(file_id)
    }

@app.websocket("/ws/progress/{file_id}")
async def websocket_endpoint(websocket: WebSocket, file_id: str):
    await websocket.accept()
    try:
        while True:
            # Send progress updates
            progress = chunk_processor.get_progress(file_id)
            await websocket.send_json({"progress": progress})
            await asyncio.sleep(1)
    except WebSocketDisconnect:
        pass

class AudioTranscriptionRequest(BaseModel):
    audio_data: str  # Base64 encoded audio data
    language_code: str = "en-US"  # Default to English

# Initialize Whisper model at startup
logger.debug("Loading Whisper model at startup")
whisper_model = whisper.load_model("base")
logger.debug("Whisper model loaded successfully")

@app.post("/api/transcribe-recording")
async def transcribe_recording(request: AudioTranscriptionRequest):
    if not whisper_model:
        raise HTTPException(status_code=503, detail="Transcription service is not available")
        
    try:
        # Add debug logging for request size
        audio_size = len(request.audio_data.encode('utf-8'))
        logger.debug(f"Received transcription request with audio data size: {audio_size / (1024 * 1024):.2f}MB")
        
        if audio_size > settings.MAX_UPLOAD_SIZE:
            logger.error(f"Audio data size ({audio_size} bytes) exceeds limit ({settings.MAX_UPLOAD_SIZE} bytes)")
            raise HTTPException(status_code=413, detail=f"Audio data size exceeds the maximum limit of {settings.MAX_UPLOAD_SIZE / (1024 * 1024):.1f}MB")
        
        # Decode base64 audio data
        try:
            audio_bytes = base64.b64decode(request.audio_data)
            logger.debug(f"Successfully decoded audio data, size: {len(audio_bytes)} bytes")
        except Exception as e:
            logger.error(f"Error decoding base64 audio data: {str(e)}")
            raise HTTPException(status_code=400, detail="Invalid audio data format")
        
        # Create a temporary file for the audio
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.webm') as temp_file:
                temp_file.write(audio_bytes)
                temp_file_path = temp_file.name
                logger.debug(f"Created temporary file: {temp_file_path}")
                logger.debug(f"Temporary file size: {os.path.getsize(temp_file_path) / (1024 * 1024):.2f}MB")
        except Exception as e:
            logger.error(f"Error creating temporary file: {str(e)}")
            raise HTTPException(status_code=500, detail="Error processing audio file")
        
        try:
            # Use the pre-loaded Whisper model
            logger.debug("Starting transcription with pre-loaded model")
            result = whisper_model.transcribe(temp_file_path)
            transcript = result["text"]
            logger.debug(f"Transcription completed, length: {len(transcript)} characters")
            
            # Generate notes from the transcript
            notes_prompt = f"""Please generate comprehensive lecture notes from this transcript. 
            Format the notes in markdown with the following structure:
            
            # Lecture Notes
            
            ## Key Points
            - Main concepts and ideas
            
            ## Detailed Notes
            - Important details and explanations
            
            ## Summary
            - Brief overview of the main topics
            
            Transcript:
            {transcript}
            """
            
            try:
                logger.debug("Generating notes")
                notes_response = llm.invoke(notes_prompt)
                notes = notes_response.content if hasattr(notes_response, 'content') else str(notes_response)
                logger.debug("Notes generation completed")
            except Exception as e:
                logger.error(f"Error generating notes: {str(e)}")
                notes = "Error generating notes. Please try again."
            
            return {
                "transcript": transcript,
                "notes": notes
            }
            
        except Exception as e:
            logger.error(f"Error during transcription: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail="Error processing audio transcription. Please try again."
            )
            
        finally:
            # Clean up the temporary file
            try:
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
                    logger.debug("Cleaned up temporary file")
            except Exception as e:
                logger.warning(f"Error removing temporary file: {str(e)}")
                
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error transcribing recording: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail="An unexpected error occurred. Please try again."
        )

@app.post("/api/summarize")
async def summarize_file(
    file: UploadFile = File(...),
    output_type: str = Form("summary")
):
    try:
        # Validate file extension
        file_extension = os.path.splitext(file.filename.lower())[1]
        settings = get_settings()
        
        logger.debug(f"Processing file: {file.filename} with extension: {file_extension}")
        
        if file_extension not in settings.ALLOWED_FILE_TYPES:
            logger.warning(f"Unsupported file type: {file_extension}. Allowed types: {settings.ALLOWED_FILE_TYPES}")
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type. Allowed types: {', '.join(settings.ALLOWED_FILE_TYPES)}"
            )
            
        # Create a temporary file to store the uploaded content
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            content = await file.read()
            temp_file.write(content)
            temp_file_path = temp_file.name
            logger.debug(f"Created temporary file at: {temp_file_path}")

        try:
            # Extract text based on file type
            logger.debug(f"Extracting text from {file_extension} file")
            
            if file_extension == '.pdf':
                text = extract_pdf_text(temp_file_path)
                logger.debug(f"Extracted {len(text)} characters from PDF")
            elif file_extension in ['.pptx', '.ppt']:
                text = extract_pptx_text(temp_file_path)
                logger.debug(f"Extracted {len(text)} characters from PowerPoint")
            elif file_extension in ['.docx', '.doc']:
                text = extract_docx_text(temp_file_path)
                logger.debug(f"Extracted {len(text)} characters from Word document")
            else:
                # For text files, read directly
                try:
                    with open(temp_file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    logger.debug(f"Extracted {len(text)} characters from text file using UTF-8")
                except UnicodeDecodeError:
                    # Try a different encoding if UTF-8 fails
                    logger.warning("UTF-8 decoding failed, trying latin-1 encoding")
                    with open(temp_file_path, 'r', encoding='latin-1') as f:
                        text = f.read()
                    logger.debug(f"Extracted {len(text)} characters from text file using latin-1")

            # Check if we got any text content
            if not text or len(text.strip()) == 0:
                logger.warning(f"No text could be extracted from {file.filename}")
                raise HTTPException(
                    status_code=400,
                    detail="Could not extract any text from the provided file."
                )

            # Generate summary or quiz based on output type
            logger.debug(f"Generating {output_type} from extracted text")
            if output_type == "quiz":
                result = generate_quiz(text)
            else:
                result = generate_summary(text)

            # Generate notes
            logger.debug("Generating notes from extracted text")
            notes = generate_notes(text)
            
            logger.debug("Successfully processed file and generated output")
            return {
                "summary": result if output_type == "summary" else None,
                "quiz": result if output_type == "quiz" else None,
                "notes": notes
            }

        finally:
            # Clean up the temporary file
            try:
                os.unlink(temp_file_path)
                logger.debug(f"Cleaned up temporary file: {temp_file_path}")
            except Exception as e:
                logger.error(f"Error cleaning up temporary file: {str(e)}")

    except HTTPException as he:
        # Re-raise HTTP exceptions
        raise he
    except Exception as e:
        logger.error(f"Error processing file: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail=f"Error processing file: {str(e)}"
        )

@app.exception_handler(StarletteHTTPException)
async def http_exception_handler(request: Request, exc: StarletteHTTPException):
    if exc.status_code == 413:  # Request Entity Too Large
        return JSONResponse(
            status_code=413,
            content={
                "detail": f"File size exceeds the maximum limit of {settings.MAX_UPLOAD_SIZE / (1024 * 1024):.1f}MB"
            }
        )
    return JSONResponse(
        status_code=exc.status_code,
        content={"detail": str(exc.detail)}
    )

class AudioChunk(BaseModel):
    chunk_data: str  # Base64 encoded audio chunk
    chunk_index: int
    total_chunks: int

class TranscriptionProgress(BaseModel):
    chunk_index: int
    total_chunks: int
    transcript: str
    is_complete: bool

# Add a queue for processing audio chunks
class AudioChunkProcessor:
    def __init__(self, max_workers=2):
        self.queue = Queue()
        self.workers = []
        self.results = {}
        self.lock = threading.Lock()
        self.progress = {}
        
        # Start worker threads
        for _ in range(max_workers):
            worker = Thread(target=self._process_chunks, daemon=True)
            worker.start()
            self.workers.append(worker)
    
    def add_chunk(self, chunk_id: str, chunk: AudioChunk, callback):
        """Add a chunk to the processing queue"""
        self.queue.put((chunk_id, chunk, callback))
        with self.lock:
            self.progress[chunk_id] = {
                'processed_chunks': 0,
                'total_chunks': chunk.total_chunks,
                'transcript': '',
                'is_complete': False
            }
    
    def _process_chunks(self):
        """Worker thread to process chunks"""
        while True:
            try:
                chunk_id, chunk, callback = self.queue.get()
                
                # Process the chunk
                try:
                    result = callback(chunk)
                    with self.lock:
                        self.results[chunk_id] = self.results.get(chunk_id, [])
                        self.results[chunk_id].append((chunk.chunk_index, result))
                        self.progress[chunk_id]['processed_chunks'] += 1
                        
                        # Check if all chunks are processed
                        if self.progress[chunk_id]['processed_chunks'] == chunk.total_chunks:
                            # Combine transcripts in order
                            sorted_results = sorted(self.results[chunk_id], key=lambda x: x[0])
                            combined_transcript = ' '.join([r[1] for r in sorted_results])
                            self.progress[chunk_id]['transcript'] = combined_transcript
                            self.progress[chunk_id]['is_complete'] = True
                except Exception as e:
                    logger.error(f"Error processing chunk {chunk_id}: {str(e)}")
                    with self.lock:
                        self.progress[chunk_id]['error'] = str(e)
                
                self.queue.task_done()
            except Exception as e:
                logger.error(f"Error in chunk processor: {str(e)}")
    
    def get_progress(self, chunk_id: str) -> TranscriptionProgress:
        """Get the progress for a chunk"""
        with self.lock:
            progress = self.progress.get(chunk_id, {})
            return TranscriptionProgress(
                chunk_index=progress.get('processed_chunks', 0),
                total_chunks=progress.get('total_chunks', 0),
                transcript=progress.get('transcript', ''),
                is_complete=progress.get('is_complete', False)
            )
    
    def get_error(self, chunk_id: str) -> str:
        """Get any error that occurred during processing"""
        with self.lock:
            return self.progress.get(chunk_id, {}).get('error')

# Create a global chunk processor
audio_chunk_processor = AudioChunkProcessor()

def process_audio_chunk(chunk: AudioChunk) -> str:
    """Process a single audio chunk with Whisper"""
    try:
        # Decode base64 audio data
        audio_bytes = base64.b64decode(chunk.chunk_data)
        
        # Create a temporary file for the chunk
        with tempfile.NamedTemporaryFile(delete=False, suffix='.webm') as temp_file:
            temp_file.write(audio_bytes)
            temp_file_path = temp_file.name
        
        try:
            # Use Whisper to transcribe the chunk
            result = whisper_model.transcribe(temp_file_path)
            return result["text"]
        finally:
            # Clean up the temporary file
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
    except Exception as e:
        logger.error(f"Error processing audio chunk: {str(e)}")
        raise

@app.post("/api/transcribe-chunk")
async def transcribe_chunk(chunk: AudioChunk):
    """Handle a single chunk of audio data"""
    try:
        chunk_id = f"chunk_{uuid.uuid4()}"
        audio_chunk_processor.add_chunk(chunk_id, chunk, process_audio_chunk)
        return {"chunk_id": chunk_id}
    except Exception as e:
        logger.error(f"Error processing chunk: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/transcription-progress/{chunk_id}")
async def get_transcription_progress(chunk_id: str):
    """Get the progress of a transcription job"""
    try:
        progress = audio_chunk_processor.get_progress(chunk_id)
        error = audio_chunk_processor.get_error(chunk_id)
        
        if error:
            raise HTTPException(status_code=500, detail=error)
            
        return progress
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error getting progress: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

class AudioFileUpload(BaseModel):
    file: UploadFile
    output_type: str = "transcript"  # Can be "transcript" or "notes"

@app.post("/api/transcribe-audio")
async def transcribe_audio_file(
    file: UploadFile = File(...),
    output_type: str = Form("transcript")
):
    if not whisper_model:
        raise HTTPException(status_code=503, detail="Transcription service is not available")
        
    try:
        # Validate file type
        if not file.content_type.startswith('audio/'):
            raise HTTPException(
                status_code=400,
                detail="Invalid file type. Please upload an audio file (MP3, WAV, M4A, etc.)"
            )
        
        # Validate file size
        file_size = 0
        content = await file.read()
        file_size = len(content)
        
        if file_size > settings.MAX_UPLOAD_SIZE:
            raise HTTPException(
                status_code=413,
                detail=f"File size exceeds the maximum limit of {settings.MAX_UPLOAD_SIZE / (1024 * 1024):.1f}MB"
            )
        
        # Create a temporary file for the audio
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file.filename.split(".")[-1]}') as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name
                logger.debug(f"Created temporary file: {temp_file_path}")
                logger.debug(f"Temporary file size: {os.path.getsize(temp_file_path) / (1024 * 1024):.2f}MB")
        except Exception as e:
            logger.error(f"Error creating temporary file: {str(e)}")
            raise HTTPException(status_code=500, detail="Error processing audio file")
        
        try:
            # Use the pre-loaded Whisper model
            logger.debug("Starting transcription with pre-loaded model")
            result = whisper_model.transcribe(temp_file_path)
            transcript = result["text"]
            logger.debug(f"Transcription completed, length: {len(transcript)} characters")
            
            response = {"transcript": transcript}
            
            # Generate notes if requested
            if output_type == "notes":
                notes_prompt = f"""Please generate comprehensive lecture notes from this transcript. 
                Format the notes in markdown with the following structure:
                
                # Lecture Notes
                
                ## Key Points
                - Main concepts and ideas
                
                ## Detailed Notes
                - Important details and explanations
                
                ## Summary
                - Brief overview of the main topics
                
                Transcript:
                {transcript}
                """
                
                try:
                    logger.debug("Generating notes")
                    notes_response = llm.invoke(notes_prompt)
                    notes = notes_response.content if hasattr(notes_response, 'content') else str(notes_response)
                    logger.debug("Notes generation completed")
                    response["notes"] = notes
                except Exception as e:
                    logger.error(f"Error generating notes: {str(e)}")
                    response["notes"] = "Error generating notes. Please try again."
            
            return response
            
        except Exception as e:
            logger.error(f"Error during transcription: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail="Error processing audio transcription. Please try again."
            )
            
        finally:
            # Clean up the temporary file
            try:
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
                    logger.debug("Cleaned up temporary file")
            except Exception as e:
                logger.warning(f"Error removing temporary file: {str(e)}")
                
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error transcribing audio file: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail="An unexpected error occurred. Please try again."
        )

@app.post("/api/response")
async def get_chat_response(request: ChatRequest) -> ChatResponse:
    try:
        # Extract content from request
        content = request.content
        user_query = request.user_query

        # Prepare the context for the LLM
        context = ""
        if content.get("transcript"):
            context += f"Transcript:\n{content['transcript']}\n\n"
        if content.get("notes"):
            context += f"Notes:\n{content['notes']}\n\n"
        if content.get("summary"):
            context += f"Summary:\n{content['summary']}\n\n"

        # Prepare the prompt for the LLM
        prompt = f"""You are a helpful AI assistant. You have access to the following content:

{context}

Please answer the following question based on the content above:
{user_query}

Remember to:
1. Only use information from the provided content
2. If the answer cannot be found in the content, say so
3. Be concise but thorough
4. Use markdown formatting for better readability
5. In the event that the user does not upload any content or file, give them a general response but politely remind them to upload a file or content and get help from the assistant.

Answer:"""

        # Get response from LLM
        response = llm.invoke(prompt).content

        return ChatResponse(response=response, user_query=user_query)
    except Exception as e:
        logger.error(f"Error in get_chat_response: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# Functions to extract text from PowerPoint and Word documents
def extract_pptx_text(file_path):
    """Extract text from PowerPoint file."""
    try:
        presentation = Presentation(file_path)
        text = []
        
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            # Join all text from this slide
            if slide_text:
                text.append(" ".join(slide_text))
        
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text from PowerPoint: {str(e)}")

def extract_docx_text(file_path):
    """Extract text from Word document."""
    try:
        doc = docx.Document(file_path)
        paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        logger.error(f"Error extracting text from Word document: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text from Word document: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    import os
    port = int(os.environ.get("PORT", 8000))
    host = "0.0.0.0"  # Important for Render
    uvicorn.run(app, host=host, port=port)

